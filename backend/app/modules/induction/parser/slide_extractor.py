from pptx.slide import Slide

def extract_slide_text(slide: Slide) -> tuple[str, str]:
    """
    Extracts the slide title and body text from a slide shape list.
    """
    title = ""
    body_elements = []

    # Try extracting title from standard title shape
    if slide.shapes.title:
        title = slide.shapes.title.text.strip()

    for shape in slide.shapes:
        # If shape has a text frame and is not the title shape
        if shape.has_text_frame and shape != slide.shapes.title:
            text = shape.text_frame.text.strip()
            if text:
                body_elements.append(text)

    # If title is still empty, try fallback to first element
    if not title and body_elements:
        first_text = body_elements[0]
        if len(first_text) < 100:
            title = first_text
            body_elements = body_elements[1:]
        else:
            title = "Untitled Slide"

    return title, "\n".join(body_elements)
