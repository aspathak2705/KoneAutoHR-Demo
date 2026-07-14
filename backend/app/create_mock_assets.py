from pptx import Presentation
from openpyxl import Workbook

def generate_mock_pptx(dest_path: str):
    prs = Presentation()

    # Slide 1: Welcome
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Welcome to KONE!"
    slide1.placeholders[1].text = "AI-Driven HR Onboarding Portal"
    slide1.notes_slide.notes_text_frame.text = "Warmly welcome all new hires. Inform them that the session is recorded for HR attendance compliance."

    # Slide 2: Core Values
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "KONE Culture & Core Values"
    slide2.placeholders[1].text = "- Safety First\n- Customer Focus\n- Collaboration & Trust"
    slide2.notes_slide.notes_text_frame.text = "Detail our commitment to safety. Explain that KONE employees are empowered to stop unsafe work."

    # Slide 3: Next Steps
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Your Next Onboarding Steps"
    slide3.placeholders[1].text = "- Setup IT Portal Profile\n- Upload Bank details\n- Complete Mandatory Safety Course"
    slide3.notes_slide.notes_text_frame.text = "IT portal credentials will arrive via email today. Complete the safety course by Friday."

    prs.save(dest_path)
    print(f"Generated mock presentation: {dest_path}")

def generate_mock_xlsx(dest_path: str):
    wb = Workbook()
    ws = wb.active
    ws.title = "New Hires"

    # Append Header with alternate column labels to verify dynamic mapping robustness
    ws.append(["Full Name", "email_address", "Function", "Role", "Office", "Start Date"])

    # Append Employees
    ws.append(["Rahul Sharma", "rahul.sharma@example.com", "Engineering", "Software Engineer", "Hyderabad", "2026-07-09"])
    ws.append(["Alice Vance", "alice.vance@example.com", "HR", "Recruiter", "Bangalore", "2026-07-09"])
    ws.append(["Bob Carter", "bob.carter@example.com", "Engineering", "DevOps Engineer", "Hyderabad", "2026-07-09"])

    wb.save(dest_path)
    wb.close()
    print(f"Generated mock employees Excel: {dest_path}")

if __name__ == "__main__":
    generate_mock_pptx("mock_presentation.pptx")
    generate_mock_xlsx("mock_employees.xlsx")
