from sqlalchemy.orm import Session as DBSession
from typing import List, Optional
from fastapi import UploadFile
from pathlib import Path
from app.repositories.presentation_repository import presentation_repository
from app.repositories.presentation_metadata_repository import presentation_metadata_repository
from app.models.presentation import Presentation
from app.services.storage_service import storage_service
from app.db.unit_of_work import UnitOfWork

class PresentationService:
    def get_all(self, db: DBSession, skip: int = 0, limit: int = 100) -> List[Presentation]:
        return presentation_repository.get_all(db, skip, limit)

    def get(self, db: DBSession, id: str) -> Optional[Presentation]:
        return presentation_repository.get(db, id)

    async def create_presentation(self, db: DBSession, name: str, file: UploadFile) -> Presentation:
        sanitized, storage_path, size = await storage_service.save_presentation_file(file)
        
        with UnitOfWork(db):
            # Create presentation record
            pres = presentation_repository.create(
                db,
                name=name,
                original_filename=file.filename,
                storage_path=storage_path
            )
            
            # Create metadata child record (defaults)
            presentation_metadata_repository.create(
                db,
                presentation_id=pres.id,
                slide_count=0,
                generation_status="PENDING"
            )
            
        # Refresh the presentation object to populate database-generated defaults
        db.refresh(pres)
        return pres

    def update(self, db: DBSession, id: str, **kwargs) -> Presentation:
        pres = presentation_repository.get(db, id)
        if not pres:
            raise ValueError(f"Presentation with id {id} not found")
        with UnitOfWork(db):
            res = presentation_repository.update(db, pres, **kwargs)
        db.refresh(res)
        return res

    def delete(self, db: DBSession, id: str) -> Optional[Presentation]:
        # Delete from disk
        pres = presentation_repository.get(db, id)
        if pres and pres.storage_path:
            import os
            try:
                os.remove(pres.storage_path)
            except OSError:
                pass
        with UnitOfWork(db):
            res = presentation_repository.delete(db, id)
        return res

presentation_service = PresentationService()


class SlideExtractor:
    def extract_slides(self, ppt_path: str, output_dir: Path) -> int:
        raise NotImplementedError


class PowerPointSlideExtractor(SlideExtractor):
    def extract_slides(self, ppt_path: str, output_dir: Path) -> int:
        """
        Automates Microsoft PowerPoint on Windows via win32com to export slide images.
        Falls back to generating dummy slide images with PIL if win32com is unavailable.
        """
        output_dir.mkdir(parents=True, exist_ok=True)
        abs_ppt = str(Path(ppt_path).resolve())
        abs_out = str(Path(output_dir).resolve())
        
        try:
            import win32com.client
            import pythoncom
            pythoncom.CoInitialize()
            import subprocess
            try:
                subprocess.run(["taskkill", "/f", "/im", "powerpnt.exe"], capture_output=True)
            except Exception:
                pass
            try:
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                # Open presentation in read-only and without a visible window if possible
                presentation = powerpoint.Presentations.Open(abs_ppt, ReadOnly=True, WithWindow=False)
                # Export slides to the folder as PNG
                presentation.Export(abs_out, "PNG")
                presentation.Close()
                powerpoint.Quit()
                del presentation
                del powerpoint
                import gc
                gc.collect()
                logger.info(f"PowerPointSlideExtractor | Successfully exported slides via win32com to {abs_out}")
            finally:
                pythoncom.CoUninitialize()
        except Exception as e:
            logger.warning(f"PowerPointSlideExtractor | win32com failure: {e}. Generating placeholders...")
            from pptx import Presentation as PptxPresentation
            prs = PptxPresentation(ppt_path)
            for idx in range(len(prs.slides)):
                slide_img = output_dir / f"slide_{idx+1}.png"
                try:
                    from PIL import Image, ImageDraw
                    img = Image.new('RGB', (1280, 720), color=(50, 50, 60))
                    d = ImageDraw.Draw(img)
                    d.text((100, 300), f"Slide {idx+1}", fill=(255, 255, 255))
                    img.save(slide_img)
                except Exception:
                    slide_img.write_bytes(b"")
            logger.info(f"PowerPointSlideExtractor | Created placeholder slide images in {abs_out}")
            
        # Standardize slide image names to slide_001.png, slide_002.png etc.
        # Resolve unique paths to prevent duplicates on case-insensitive filesystems
        raw_files = list(output_dir.glob("*.png")) + list(output_dir.glob("*.PNG"))
        unique_paths = {}
        for f in raw_files:
            try:
                unique_paths[f.resolve()] = f
            except Exception:
                unique_paths[f] = f
        exported_files = list(unique_paths.values())
        import re
        slides_count = 0
        for f in exported_files:
            match = re.search(r'\d+', f.name)
            if match:
                num = int(match.group())
                new_name = output_dir / f"slide_{num:03d}.png"
                if f != new_name:
                    if new_name.exists():
                        new_name.unlink()
                    f.rename(new_name)
                slides_count = max(slides_count, num)
        return slides_count


class SlideFingerprintGenerator:
    def generate_fingerprints(self, slides_dir: Path, output_file: Path) -> dict:
        """
        Computes perceptual hashes and slide dimensions metadata for fingerprints.json.
        """
        fingerprints = {}
        slide_files = sorted(list(slides_dir.glob("slide_*.png")))
        
        for sf in slide_files:
            slide_name = sf.name
            phash, width, height, aspect_ratio = self._analyze_image(sf)
            fingerprints[slide_name] = {
                "phash": phash,
                "thumbnail": f"presentation_assets/slides/{slide_name}",
                "width": width,
                "height": height,
                "aspect_ratio": aspect_ratio,
                "edge_descriptor": [0.1, 0.2, 0.1, 0.3]
            }
            
        with open(output_file, "w", encoding="utf-8") as f:
            json.dump(fingerprints, f, indent=2)
            
        logger.info(f"SlideFingerprintGenerator | Saved {len(fingerprints)} slide fingerprints to {output_file}")
        return fingerprints

    def _analyze_image(self, img_path: Path):
        try:
            from PIL import Image
            img = Image.open(img_path)
            width, height = img.size
            aspect_ratio = width / height if height > 0 else 1.0
            
            gray = img.convert('L').resize((8, 8), Image.Resampling.LANCZOS)
            pixels = list(gray.getdata())
            avg = sum(pixels) / 64.0
            phash = "".join("1" if p >= avg else "0" for p in pixels)
            
            return phash, width, height, aspect_ratio
        except Exception as e:
            logger.error(f"SlideFingerprintGenerator | Failed to analyze {img_path}: {e}")
            return "0" * 64, 1280, 720, 1.777


class PresentationAssetBuilder:
    def build_assets(self, session_id: str, ppt_path: str) -> None:
        """
        Main orchestration class to build presentation assets.
        """
        from app.services.storage_service import storage_service
        session_dir = storage_service.get_session_dir(session_id)
        assets_dir = session_dir / "presentation_assets"
        slides_dir = assets_dir / "slides"
        
        slides_dir.mkdir(parents=True, exist_ok=True)
        
        # 1. Extract slides using PowerPointSlideExtractor
        extractor = PowerPointSlideExtractor()
        slide_count = extractor.extract_slides(ppt_path, slides_dir)
        
        # 2. Create deterministic presentation timeline metadata.
        self.generate_presentation_manifest(session_id, ppt_path, slide_count, assets_dir)
        
        # 3. Create metadata.json
        metadata = {
            "session_id": session_id,
            "ppt_path": ppt_path,
            "slide_count": slide_count,
            "generated_at": datetime.datetime.now().isoformat()
        }
        with open(assets_dir / "metadata.json", "w", encoding="utf-8") as f:
            json.dump(metadata, f, indent=2)
            
        logger.info(f"PresentationAssetBuilder | Completed building assets for session {session_id}")

    def generate_presentation_manifest(
        self,
        session_id: str,
        ppt_path: str,
        slide_count: int,
        assets_dir: Path,
    ) -> dict:
        from app.modules.meeting_bot.media.audio_controller import get_audio_controller

        audio_ctrl = get_audio_controller(session_id)
        slides = []
        for slide_num in range(1, slide_count + 1):
            audio_file = f"slide_{slide_num}.mp3"
            slides.append({
                "number": slide_num,
                "audio": audio_file,
                "duration": audio_ctrl.get_duration(audio_file),
            })

        manifest = {
            "session_id": session_id,
            "ppt_path": ppt_path,
            "slides": slides,
            "generated_at": datetime.datetime.now().isoformat(),
        }
        with open(assets_dir / "presentation.json", "w", encoding="utf-8") as f:
            json.dump(manifest, f, indent=2)
        logger.info(
            f"PresentationAssetBuilder | Saved deterministic presentation manifest for {len(slides)} slides"
        )
        return manifest

# Additional missing imports for newly added classes
import json
import datetime
from pathlib import Path
from loguru import logger



