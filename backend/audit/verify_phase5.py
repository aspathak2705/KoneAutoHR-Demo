import asyncio
import os
import json
import time
from pathlib import Path
from sqlalchemy.orm import Session as DBSession

from app.db.database import SessionLocal
from app.models.session import Session
from app.models.meeting import Meeting
from app.models.presentation import Presentation
from app.models.presentation_script import PresentationScript
from app.models.presentation_question import PresentationQuestion
from app.models.employee_list import EmployeeList
from app.models.organization_config import OrganizationConfig

# Import missing dependency models to register correctly
from app.models.presentation_asset import PresentationAsset
from app.models.presentation_job import PresentationJob
from app.models.upload import Upload
from app.models.runtime import Runtime
from app.models.runtime_message import RuntimeMessage
from app.models.attendance import Attendance
from app.models.invitation_draft import InvitationDraft

from app.modules.presentation_observer.models.observation import Observation
from app.modules.presentation_observer.models.observation_state import ObservationState
from app.modules.presentation_observer.models.observation_event import ObservationEvent
from app.modules.semantic_browser.models.presentation_state import PresentationMode
from app.modules.semantic_browser.models.meeting_state import MeetingState

from app.modules.induction_runtime.services.induction_runtime_service import induction_runtime_service
from app.modules.induction_runtime.models.runtime_state import RuntimeState
from app.modules.induction_runtime.config.runtime_config import RuntimeConfig
from app.modules.induction_runtime.orchestrator.runtime_coordinator import RuntimeCoordinator
from app.modules.induction_runtime.narrator.voice_output_interface import VoiceOutputInterface

class TracedVoiceOutput(VoiceOutputInterface):
    def __init__(self):
        self.history = []
        self.is_speaking = False
        self._current_callback = None

    def say(self, text: str, callback = None, **kwargs) -> None:
        self.is_speaking = True
        self.history.append(("say", text))
        self._current_callback = callback
        
    def interrupt(self) -> None:
        self.is_speaking = False
        self.history.append(("interrupt", None))
        
    def stop(self) -> None:
        self.is_speaking = False
        self.history.append(("stop", None))
        
    def resume(self) -> None:
        self.is_speaking = True
        self.history.append(("resume", None))

    def trigger_callback(self):
        self.is_speaking = False
        if self._current_callback:
            cb = self._current_callback
            self._current_callback = None
            if asyncio.iscoroutinefunction(cb):
                asyncio.create_task(cb())
            else:
                cb()

async def run_verification():
    """
    Phase 5 Verification: Runtime Engine & Conductor.
    """
    assertions = 0
    start_time = asyncio.get_event_loop().time()
    warnings = []
    
    db = SessionLocal()
    session_id = "test-session-audit-p5"
    excel_path = "mock_inductees_audit.xlsx"

    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.append(["Name", "Email", "Department", "Designation", "Location", "Joining Date"])
    ws.append(["Jane Doe", "jane.doe@kone.com", "HR", "Talent Acquisition Coordinator", "Espoo", "2026-08-01"])
    wb.save(excel_path)

    from pptx import Presentation as PPTXPresentation
    prs = PPTXPresentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save("mock_deck.pptx")

    def clean_db(db_session, s_id):
        try:
            db_session.query(Meeting).filter(Meeting.session_id == s_id).delete()
            db_session.query(Session).filter(Session.id == s_id).delete()
            db_session.query(PresentationScript).filter(PresentationScript.presentation_id == "pres-audit").delete()
            db_session.query(PresentationQuestion).filter(PresentationQuestion.presentation_id == "pres-audit").delete()
            db_session.query(Presentation).filter(Presentation.id == "pres-audit").delete()
            db_session.query(EmployeeList).filter(EmployeeList.id == "emplist-audit").delete()
            db_session.commit()
        except Exception as e:
            db_session.rollback()
            warnings.append(f"DB cleanup warning: {e}")

    try:
        clean_db(db, session_id)

        # Seed Config only if none exists
        if not db.query(OrganizationConfig).first():
            config = OrganizationConfig(
                company_name="KONE",
                company_domain="kone.com",
                ai_officer_name="KONE AI Officer",
                ai_trainer_name="KONE Trainer",
                ai_role_description="AI Onboarding Assistant",
                vocal_tone="Professional",
                communication_style="Direct"
            )
            db.add(config)
            db.commit()

        # Seed Presentation
        pres = Presentation(
            id="pres-audit",
            name="KONE Onboarding deck",
            original_filename="kone_onboarding.pptx",
            storage_path="mock_deck.pptx"
        )
        db.add(pres)

        # Seed EmployeeList
        emp_list = EmployeeList(
            id="emplist-audit",
            name="Inductees August 2026",
            original_filename="inductees.xlsx",
            storage_path=excel_path,
            employee_count=1
        )
        db.add(emp_list)

        # Seed Session
        sess = Session(
            id=session_id,
            name="Jane Doe Induction Session",
            presentation_id="pres-audit",
            employee_list_id="emplist-audit",
            status="PENDING"
        )
        db.add(sess)

        # Seed Meeting
        meeting = Meeting(
            session_id=session_id,
            teams_meeting_url="https://teams.microsoft.com/l/meetup-join/test",
            meeting_passcode="123456",
            organizer_name="HR Team",
            meeting_date="2026-08-01",
            meeting_time="10:00 AM"
        )
        db.add(meeting)

        # Seed PresentationScript
        script_payload = {
            "welcome_flow": {
                "greeting": "Hello and welcome to KONE onboarding.",
                "presenter_intro": "I am your AI HR Trainer {trainer_name}, here to guide you today.",
                "session_rules": "Please stay muted during slides.",
                "agenda": "Today we cover KONE values and safety guidelines."
            },
            "slide_narrations": {
                "1": "First slide: KONE is built on innovation and safety. Welcome {trainer_name} from {company_name}.",
                "2": "Second slide: Safety is our core priority. Report hazards immediately."
            },
            "closing_script": "That concludes KONE values and safety rules induction."
        }
        script = PresentationScript(
            presentation_id="pres-audit",
            script_content=json.dumps(script_payload),
            llm_model="gemini-1.5-pro",
            status="ACTIVE"
        )
        db.add(script)

        # Seed PresentationQuestion
        faq_payload = [
            {"question": "What is the policy on safety?", "answer": "KONE has a zero-tolerance policy for safety violations."}
        ]
        faq = PresentationQuestion(
            presentation_id="pres-audit",
            questions_content=json.dumps(faq_payload),
            status="ACTIVE"
        )
        db.add(faq)
        db.commit()

        # Initialize Coordinator
        traced_voice = TracedVoiceOutput()
        coord = RuntimeCoordinator(db, session_id, voice_output=traced_voice)
        coord.initialize()
        
        assert coord.session_manager.state == RuntimeState.CREATED
        assertions += 1

        # Simulate Lobby first
        obs_lobby = Observation(
            timestamp=time.time(),
            meeting_state=MeetingState.LOBBY,
            observation_state=ObservationState.WAITING,
            presentation_state=PresentationMode.WAITING_SCREEN,
            events=[],
            timeline_index=0
        )
        await coord.poll_cycle(obs_lobby)
        assert coord.session_manager.state == RuntimeState.WAITING_FOR_PRESENTATION
        assertions += 1

        # Simulate Started
        obs_started = Observation(
            timestamp=time.time(),
            meeting_state=MeetingState.CONNECTED,
            observation_state=ObservationState.ACTIVE,
            presentation_state=PresentationMode.POWERPOINT_SHARED,
            events=[ObservationEvent.PRESENTATION_STARTED],
            timeline_index=1
        )
        await coord.poll_cycle(obs_started)
        await asyncio.sleep(0.1)
        assert coord.session_manager.state == RuntimeState.INTRODUCTION
        assertions += 1
        
        # Complete Welcome
        traced_voice.trigger_callback()
        await asyncio.sleep(0.1)
        assert coord.session_manager.state == RuntimeState.PRESENTING
        assertions += 1

        # Slide 2 change
        obs_slide2 = Observation(
            timestamp=time.time(),
            meeting_state=MeetingState.CONNECTED,
            observation_state=ObservationState.ACTIVE,
            presentation_state=PresentationMode.POWERPOINT_SHARED,
            events=[ObservationEvent.SLIDE_CHANGED],
            timeline_index=2
        )
        await coord.poll_cycle(obs_slide2)
        await asyncio.sleep(0.1)
        assert coord.memory.current_slide_number == 2
        assertions += 1

        # Q&A Inject
        ans = await coord.inject_question("Jane Doe", "What is the policy on safety?")
        assert "zero-tolerance" in ans
        assertions += 1

        # Stage 7: Pipeline Decoupling & Fast Failure
        from app.modules.induction.services.preparation_orchestrator import preparation_orchestrator
        
        # Test 1: Fast failure when script is not generated
        session_decoupled_id = "test-session-decoupled-audit"
        
        # Create temp session without script
        sess_decoupled = Session(
            id=session_decoupled_id,
            name="Decoupled Pipeline Session Audit",
            presentation_id="pres-audit",
            employee_list_id="emplist-audit",
            status="PENDING"
        )
        db.add(sess_decoupled)
        db.commit()
        
        try:
            # Recreate jobs but SCRIPT job status is PENDING (not COMPLETED)
            script_job = preparation_orchestrator._get_or_create_job(db, session_decoupled_id, "SCRIPT")
            from app.services.presentation_job_service import presentation_job_service
            presentation_job_service.update_job_status(db, script_job.id, status="PENDING")
            # Deleting script if it exists in DB just in case
            db.query(PresentationScript).filter(PresentationScript.presentation_id == "pres-audit").delete()
            db.commit()
            
            try:
                await preparation_orchestrator.run_audio_generation(session_decoupled_id, "dummy-job-id")
                assert False, "Should have raised ValueError because SCRIPT job is not COMPLETED"
            except ValueError as val_err:
                assert "Presentation script has not been generated. Please complete script generation first." in str(val_err)
                assertions += 1
                
            # Test 2: Run preparation with mocked LLM generation
            import app.modules.induction.services.script_pipeline as sp_mod
            async def mock_generate_scripts(*args, **kwargs):
                return {
                    "welcome_flow": {
                        "greeting": "Hello and welcome to KONE onboarding.",
                        "summary": "Please stay muted during slides."
                    },
                    "slide_narrations": {
                        "1": {"narration": "First slide: KONE is built on innovation and safety. Welcome."},
                        "2": {"narration": "Second slide: Safety is our core priority. Report hazards immediately."}
                    },
                    "faq": [
                        {"question": "What is the policy on safety?", "answer": "KONE has a zero-tolerance policy for safety violations."}
                    ],
                    "closing_script": {"summary": "That concludes KONE values and safety rules induction."}
                }
            original_gen = sp_mod.generate_induction_package_scripts
            sp_mod.generate_induction_package_scripts = mock_generate_scripts
            
            try:
                await preparation_orchestrator.run_preparation(session_decoupled_id, "dummy-prep-job-id")
            finally:
                sp_mod.generate_induction_package_scripts = original_gen
            
            # Verify script is ready in DB
            db_script = db.query(PresentationScript).filter(PresentationScript.presentation_id == "pres-audit").first()
            assert db_script is not None
            assertions += 1
            
            # Verify SCRIPT job is completed
            script_job = presentation_job_service.get_job_by_session(db, session_id=session_decoupled_id, job_type="SCRIPT")
            assert script_job.status == "COMPLETED"
            assertions += 1
            
            # Test 3: Run audio generation
            await preparation_orchestrator.run_audio_generation(session_decoupled_id, "dummy-audio-job-id")
            audio_job = presentation_job_service.get_job_by_session(db, session_id=session_decoupled_id, job_type="AUDIO")
            assert audio_job.status == "COMPLETED"
            assertions += 1
            
            # Test 4: Run package generation
            await preparation_orchestrator.run_package_generation(session_decoupled_id)
            pack_job = presentation_job_service.get_job_by_session(db, session_id=session_decoupled_id, job_type="PACKAGE")
            assert pack_job.status == "COMPLETED"
            assertions += 1
            
            sess_check = db.query(Session).filter(Session.id == session_decoupled_id).first()
            assert sess_check.status == "READY"
            assertions += 1
            
        finally:
            # Cleanup decoupled session details
            db.query(Meeting).filter(Meeting.session_id == session_decoupled_id).delete()
            db.query(PresentationJob).filter(PresentationJob.session_id == session_decoupled_id).delete()
            db.query(Session).filter(Session.id == session_decoupled_id).delete()
            db.commit()

    finally:
        if os.path.exists(excel_path):
            os.remove(excel_path)
        if os.path.exists("mock_deck.pptx"):
            os.remove("mock_deck.pptx")
        clean_db(db, session_id)
        db.close()
        
    duration = (asyncio.get_event_loop().time() - start_time) * 1000 # ms
    return {
        "success": True,
        "assertions": assertions,
        "duration_ms": duration,
        "warnings": warnings
    }
