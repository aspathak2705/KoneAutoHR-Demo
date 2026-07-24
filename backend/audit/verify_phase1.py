import asyncio
import tempfile
from pathlib import Path
from pptx import Presentation
from app.modules.induction.parser.ppt_parser import parse_presentation

async def run_verification():
    """
    Phase 1 Verification: PPTX Parsing & Text Extraction.
    """
    assertions = 0
    start_time = asyncio.get_event_loop().time()
    
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        ppt_path = temp_path / "test_deck.pptx"
        
        # 1. Create a mock PPTX deck using python-pptx
        prs = Presentation()
        
        # Slide 1
        slide_layout = prs.slide_layouts[0] # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "KONE Induction Slide Title 1"
        subtitle.text = "KONE Onboarding Safety Rules Content"
        
        # Slide 2
        slide_layout = prs.slide_layouts[1] # Content Slide
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "KONE Innovation"
        content_placeholder = slide.placeholders[1]
        content_placeholder.text = "Innovation is one of our key pillars."
        
        prs.save(str(ppt_path))
        
        # 2. Parse the presentation
        slides_data = parse_presentation(str(ppt_path), temp_path)
        
        # 3. Assertions
        assert len(slides_data) == 2, "Expected exactly 2 slides parsed"
        assertions += 1
        
        assert slides_data[0]["slide_number"] == 1
        assertions += 1
        
        assert slides_data[0]["title"] == "KONE Induction Slide Title 1"
        assertions += 1
        
        assert "Safety Rules" in slides_data[0]["content"]
        assertions += 1
        
        assert slides_data[1]["title"] == "KONE Innovation"
        assertions += 1
        
        assert "pillar" in slides_data[1]["content"]
        assertions += 1
        
    duration = (asyncio.get_event_loop().time() - start_time) * 1000 # ms
    return {
        "success": True,
        "assertions": assertions,
        "duration_ms": duration,
        "warnings": []
    }
